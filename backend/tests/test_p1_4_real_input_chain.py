"""P1-4 C0 evidence: real-file import chain and restart re-read.

The chain under test is: import -> parse -> text -> span -> index -> retrieval
-> Q&A citation -> citation offset back into the stored body, followed by a
process-equivalent restart (a new application over the same data root).

Fixtures are built at runtime from real document containers (python-docx,
python-pptx) and, for PDF, from a browser-rendered multi-page document.  No
fixture is committed to the repository and no live data root is touched.
"""

from __future__ import annotations

import io
import os
import subprocess
import sys
from pathlib import Path

import pytest
from docx import Document
from docx.shared import Inches
from fastapi.testclient import TestClient
from PIL import Image
from pptx import Presentation
from pptx.util import Inches as PptxInches

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from app.config import AppConfig
from app.main import create_app

QUERY = "citation retrieval stability"
CHINESE_NAME = "中文资料-用于验证解析检索与引用定位的长文件名样本.txt"


def _png_bytes() -> bytes:
    buffer = io.BytesIO()
    Image.new("RGB", (120, 90), (215, 225, 240)).save(buffer, format="PNG")
    return buffer.getvalue()


def _client(root: Path, *, max_upload_bytes: int = 8 * 1024 * 1024) -> TestClient:
    return TestClient(create_app(AppConfig(
        data_root=root, max_upload_bytes=max_upload_bytes, ai_provider_id="fake",
    )))


def _chromium() -> Path | None:
    base = os.environ.get("LOCALAPPDATA")
    if not base:
        return None
    for candidate in sorted(Path(base).glob("ms-playwright/chromium*/chrome-win*/chrome.exe")):
        if candidate.is_file():
            return candidate
    return None


def _render_pdf(html: str, target: Path, *, assets: dict[str, bytes] | None = None) -> None:
    browser = _chromium()
    if browser is None:
        pytest.skip("browser-rendered PDF fixture requires the managed Chromium binary")
    for name, payload in (assets or {}).items():
        (target.parent / name).write_bytes(payload)  # noqa: PERF203 - fixture assets are tiny
    source = target.parent / f"{target.stem}.source.html"
    source.write_text(html, encoding="utf-8")
    subprocess.run(
        [str(browser), "--headless", "--disable-gpu", "--no-sandbox", "--no-pdf-header-footer",
         f"--print-to-pdf={target}", source.resolve().as_uri()],
        check=True, capture_output=True, timeout=180,
    )
    source.unlink()


def _text_pdf(directory: Path) -> Path:
    target = directory / "real-handbook.pdf"
    _render_pdf(
        "<html><head><meta charset='utf-8'><style>body{font-family:serif}"
        ".cols{column-count:2}.page{page-break-after:always}</style></head><body>"
        "<div class='page'><h1>Study Handbook</h1><p>Table of contents</p>"
        "<ol><li>Chapter 1 Stability</li><li>Chapter 2 Retrieval</li></ol><p>page 1</p></div>"
        "<div class='page'><h2>Chapter 1 Stability</h2><div class='cols'>"
        "<p>Verified stability keeps predictable behaviour under disturbance.</p>"
        "<p>The second column continues the same chapter with supporting sentences.</p>"
        "</div><p>page 2</p></div>"
        "<div><h2>Chapter 2 Retrieval</h2><p>Hybrid citation retrieval stability keeps every "
        "answer traceable to the original text span.</p><p>page 3</p></div></body></html>",
        target,
    )
    return target


def _image_pdf(directory: Path) -> Path:
    target = directory / "scanned-page.pdf"
    _render_pdf(
        "<html><head><meta charset='utf-8'></head><body>"
        "<img src='scan.png' style='width:420px'><img src='scan.png' style='width:420px'>"
        "</body></html>",
        target, assets={"scan.png": _png_bytes()},
    )
    return target


def _docx(directory: Path) -> Path:
    target = directory / "real-notes.docx"
    document = Document()
    document.add_heading("Study notes with citation retrieval stability", 0)
    document.add_paragraph("第一段：可验证稳定性要求每个断言都能追溯到观察结果。")
    document.add_paragraph("列表项一", style="List Bullet")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "concept"
    table.cell(0, 1).text = "citation retrieval stability"
    table.cell(1, 0).text = "rule"
    table.cell(1, 1).text = "引用必须能定位回原文"
    document.add_paragraph("第二段：retrieval 与 citation 的关系。")
    document.add_picture(str(_png(directory)), width=Inches(1))
    document.save(target)
    return target


def _pptx(directory: Path) -> Path:
    target = directory / "real-deck.pptx"
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[1])
    first.shapes.title.text = "第一页 学习节奏"
    first.placeholders[1].text = "citation retrieval stability keeps today's plan explainable"
    second = presentation.slides.add_slide(presentation.slide_layouts[1])
    second.shapes.title.text = "第二页 引用与检索"
    second.placeholders[1].text = "引用必须定位回原文 span"
    picture_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    picture_slide.shapes.add_picture(str(_png(directory)), PptxInches(1), PptxInches(1), PptxInches(3))
    presentation.save(target)
    return target


def _png(directory: Path) -> Path:
    target = directory / "embedded.png"
    if not target.exists():
        target.write_bytes(_png_bytes())
    return target


def _upload(client: TestClient, path: Path, *, name: str | None = None) -> dict[str, object]:
    response = client.post("/api/materials", files={
        "file": (name or path.name, path.read_bytes(), "application/octet-stream"),
    })
    assert response.status_code == 201, response.text
    return response.json()


def _chain(client: TestClient, material_id: str, query: str) -> dict[str, object]:
    index = client.post(f"/api/materials/{material_id}/ai-index")
    assert index.status_code == 200, index.text
    detail = client.get(f"/api/materials/{material_id}")
    assert detail.status_code == 200
    retrieval = client.post("/api/retrieval", json={
        "query": query, "material_ids": [material_id], "mode": "lexical", "top_k": 5,
    })
    assert retrieval.status_code == 200, retrieval.text
    answer = client.post("/api/qa/ask", json={
        "question": query, "material_ids": [material_id], "retrieval_mode": "lexical", "top_k": 5,
    })
    assert answer.status_code == 200, answer.text
    return {"index": index.json(), "detail": detail.json(),
            "retrieval": retrieval.json(), "answer": answer.json()}


def _assert_citation_locates_body(client: TestClient, detail: dict[str, object], answer: dict[str, object]) -> str:
    citations = answer["citations"]
    assert citations, answer
    key = str(citations[0]["citation_key"])
    located = client.get(f"/api/qa/citations/{key}")
    assert located.status_code == 200
    payload = located.json()
    assert payload["status"] == "valid"
    assert payload["material_id"] == detail["id"]
    body = str(detail["text"])
    start, end = int(payload["start_offset"]), int(payload["end_offset"])
    assert 0 <= start < end <= len(body)
    assert " ".join(body[start:end].split())[:240] == payload["excerpt"]
    assert "stored_path" not in payload
    return key


@pytest.mark.parametrize("builder,span_kind,expected_spans,query", [
    (_text_pdf, "page", 3, QUERY),
    (_docx, "document", 1, QUERY),
    (_pptx, "slide", 3, QUERY),
])
def test_real_document_chain_produces_locatable_citations(tmp_path: Path, builder, span_kind, expected_spans, query):
    source = builder(_fixture_dir(tmp_path))
    root = tmp_path / "data"
    with _client(root) as client:
        created = _upload(client, source)
        assert created["status"] == "success", created
        material_id = str(created["material_id"])
        result = _chain(client, material_id, query)
        detail = result["detail"]
        assert [span["span_kind"] for span in detail["spans"]] == [span_kind] * expected_spans
        assert [span["ordinal"] for span in detail["spans"]] == list(range(1, expected_spans + 1))
        assert len(str(detail["text"])) > 0
        assert result["index"]["status"] == "ready"
        assert int(result["index"]["chunk_count"]) >= 1
        assert result["retrieval"]["status"] == "succeeded"
        assert result["retrieval"]["hits"]
        assert result["answer"]["status"] == "succeeded"
        key = _assert_citation_locates_body(client, detail, result["answer"])

    with _client(root) as restarted:
        reopened = restarted.get(f"/api/materials/{material_id}")
        assert reopened.status_code == 200
        assert reopened.json()["text"] == detail["text"]
        assert len(reopened.json()["spans"]) == expected_spans
        status = restarted.get(f"/api/materials/{material_id}/ai-index").json()
        assert status["status"] == "ready"
        assert restarted.get(f"/api/qa/citations/{key}").json()["status"] == "valid"
        threads = restarted.get("/api/qa/threads").json()
        assert threads["items"]


def _fixture_dir(tmp_path: Path) -> Path:
    directory = tmp_path / "fixtures"
    directory.mkdir(parents=True, exist_ok=True)
    return directory


def test_real_text_and_markdown_chain_supports_chinese_names(tmp_path: Path):
    directory = _fixture_dir(tmp_path)
    chinese = directory / CHINESE_NAME
    chinese.write_text("第一节 可验证稳定性。\n第二节 检索与引用必须能定位回原文。\n", encoding="utf-8")
    markdown = directory / "real-guide.md"
    markdown.write_text("# 学习指南\n\n- 引用必须能定位回原文\n- 检索必须可复现\n", encoding="utf-8")
    root = tmp_path / "data"
    with _client(root) as client:
        for source, query in ((chinese, "检索"), (markdown, "引用")):
            created = _upload(client, source)
            assert created["status"] == "success", created
            assert created["original_name"] == source.name
            material_id = str(created["material_id"])
            result = _chain(client, material_id, query)
            assert [span["span_kind"] for span in result["detail"]["spans"]] == ["document"]
            assert result["retrieval"]["status"] == "succeeded"
            assert result["retrieval"]["hits"]
            _assert_citation_locates_body(client, result["detail"], result["answer"])


def test_image_only_pdf_reports_empty_text_without_ocr(tmp_path: Path):
    source = _image_pdf(_fixture_dir(tmp_path))
    with _client(tmp_path / "data") as client:
        created = _upload(client, source)
        assert created["status"] == "empty"
        assert created["error_code"] is None
        assert created["warnings"]
        assert all("OCR" in warning or "文字层" in warning for warning in created["warnings"])
        material_id = str(created["material_id"])
        index = client.post(f"/api/materials/{material_id}/ai-index")
        assert index.status_code == 200
        assert index.json()["status"] == "empty"
        assert index.json()["chunk_count"] == 0
        retrieval = client.post("/api/retrieval", json={
            "query": QUERY, "material_ids": [material_id], "mode": "lexical", "top_k": 5,
        })
        assert retrieval.json()["status"] == "failed"
        assert retrieval.json()["error_code"] == "retrieval_not_ready"


def test_declared_but_unsupported_extensions_report_stable_codes(tmp_path: Path):
    directory = _fixture_dir(tmp_path)
    cases = {
        "legacy.doc": (b"\xd0\xcf\x11\xe0legacy", "rejected", "requires_converter"),
        "legacy.ppt": (b"\xd0\xcf\x11\xe0legacy ppt", "rejected", "requires_converter"),
        "note.rtf": (rb"{\rtf1 hello}", "rejected", "unsupported_rtf"),
        "data.xml": (b"<r><a>1</a></r>", "rejected", "unsupported_format"),
        "broken.pdf": (b"%PDF-1.4 broken payload", "failed", "corrupt_pdf"),
        "blank.txt": (b"", "empty", None),
    }
    with _client(tmp_path / "data") as client:
        for name, (payload, status, code) in cases.items():
            source = directory / name
            source.write_bytes(payload)
            created = _upload(client, source)
            assert created["status"] == status, (name, created)
            assert created["error_code"] == code, (name, created)
            assert created["text_length"] == 0
            body = str(created)
            assert "Traceback" not in body and "SELECT" not in body and "stored_path" not in body


def test_encrypted_pdf_is_rejected_as_unreadable_without_leaking_internals(tmp_path: Path):
    from pypdf import PdfWriter

    directory = _fixture_dir(tmp_path)
    plain = _text_pdf(directory)
    encrypted = directory / "encrypted.pdf"
    writer = PdfWriter()
    writer.append(str(plain))
    writer.encrypt("p1-4-secret")
    with encrypted.open("wb") as handle:
        writer.write(handle)
    with _client(tmp_path / "data") as client:
        created = _upload(client, encrypted)
        # Current truth: an encrypted PDF is reported through the generic
        # unreadable-PDF code.  It is refused, never silently indexed.
        assert created["status"] == "failed"
        assert created["error_code"] == "corrupt_pdf"
        assert created["text_length"] == 0
        assert "p1-4-secret" not in str(created)


def test_upload_boundaries_reject_oversize_and_invalid_names(tmp_path: Path):
    with _client(tmp_path / "data", max_upload_bytes=1024) as client:
        oversize = client.post("/api/materials", files={"file": ("big.txt", b"x" * 4096, "text/plain")})
        assert oversize.status_code == 413
        assert oversize.json()["detail"] == "file_too_large"
        invalid = client.post("/api/materials", files={"file": ("bad/name.txt", b"ok", "text/plain")})
        assert invalid.status_code == 400
        assert invalid.json()["detail"] == "invalid_filename"
        assert client.get("/api/materials").json() == []


def test_shared_hash_import_keeps_one_original_and_survives_restart(tmp_path: Path):
    directory = _fixture_dir(tmp_path)
    first = directory / "same-content-one.md"
    first.write_text("# 共享内容\n\n引用必须能定位回原文。\n", encoding="utf-8")
    second = directory / "same-content-two.md"
    second.write_bytes(first.read_bytes())
    root = tmp_path / "data"
    with _client(root) as client:
        one = _upload(client, first)
        two = _upload(client, second)
        assert one["source_sha256"] == two["source_sha256"]
        assert one["material_id"] != two["material_id"]
    originals = [path for path in (root / "originals").rglob("original") if path.is_file()]
    assert len(originals) == 1
    with _client(root) as restarted:
        names = {item["original_name"] for item in restarted.get("/api/materials").json()}
        assert names == {first.name, second.name}
        for material in restarted.get("/api/materials").json():
            download = restarted.get(f"/api/materials/{material['id']}/original")
            assert download.status_code == 200
            assert download.content == first.read_bytes()


def test_shared_hash_both_materials_can_be_indexed_after_fix(tmp_path: Path):
    """After P14-P0-05 fix (migration v14), both materials can be indexed.

    Two materials with identical content are accepted by import (shared-hash
    storage is intentional). After the revision fingerprint fix (migration v14),
    the fingerprint includes material_id, so both materials can build their own
    AI index and be used independently for retrieval, Q&A, and citations.
    """
    directory = _fixture_dir(tmp_path)
    first = directory / "duplicate-one.md"
    first.write_text("# 共享内容\n\ncitation retrieval stability 必须可追溯。\n", encoding="utf-8")
    second = directory / "duplicate-two.md"
    second.write_bytes(first.read_bytes())
    root = tmp_path / "data"
    with _client(root) as client:
        one = _upload(client, first)
        two = _upload(client, second)
        # Both materials can now be indexed successfully
        assert client.post(f"/api/materials/{one['material_id']}/ai-index").status_code == 200
        assert client.post(f"/api/materials/{two['material_id']}/ai-index").status_code == 200
        # Both show ready status
        assert client.get(f"/api/materials/{one['material_id']}/ai-index").json()["status"] == "ready"
        assert client.get(f"/api/materials/{two['material_id']}/ai-index").json()["status"] == "ready"
        # Both can answer questions
        answer_one = client.post("/api/qa/ask", json={
            "question": QUERY, "material_ids": [one["material_id"]],
            "retrieval_mode": "lexical", "top_k": 5,
        })
        assert answer_one.status_code == 200
        answer_two = client.post("/api/qa/ask", json={
            "question": QUERY, "material_ids": [two["material_id"]],
            "retrieval_mode": "lexical", "top_k": 5,
        })
        assert answer_two.status_code == 200
        # Deleting one material does not affect the other
        assert client.delete(f"/api/materials/{one['material_id']}").status_code == 204
        assert client.get(f"/api/materials/{two['material_id']}/ai-index").json()["status"] == "ready"

    with _client(root) as restarted:
        assert restarted.get(f"/api/materials/{two['material_id']}/ai-index").json()["status"] == "ready"


def test_revision_fingerprint_conflict_error_mapping_exists():
    """C1 added revision_fingerprint_conflict to the error map.

    Even though P14-P0-05 fix prevents this conflict in normal use (materials
    with shared content now get distinct fingerprints), the error code mapping
    must remain for backward compatibility and edge cases.
    """
    shared = (ROOT / "app" / "static" / "js" / "api.js").read_text(encoding="utf-8")
    assert "request_failed:'请求失败，请重试'" in shared
    # C1 fix: revision_fingerprint_conflict now has a user-facing message
    assert "revision_fingerprint_conflict:'内容指纹冲突" in shared
